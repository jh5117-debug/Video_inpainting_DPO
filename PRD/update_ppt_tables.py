from copy import deepcopy
import os
from pathlib import Path
import re

from pptx import Presentation
from pptx.util import Pt


PPT_PATH = os.environ.get(
    "PPT_PATH",
    str(Path(__file__).resolve().parent / "PPT.pptx"),
)


TABLES = {
    "Native-Best-24G / BR": {
        "kind": "BR",
        "conclusion": "结论：这是 24GB 下保留方法 native-style setting 的 BR 表；ProPainter 最强，MiniMax 第二，但仍混有原生分辨率与长窗口差异。",
        "rows": [
            ["ProPainter", "35.5167", "0.9813", "0.0141", "7.4734", "0.0818", "0.9703", "0.0295", "1.2362", "69"],
            ["MiniMax-Remover", "33.1124", "0.9712", "0.0192", "3.9526", "0.0904", "0.9694", "N/A", "1.2379", "65"],
            ["DiffuEraser", "31.1135", "0.9642", "0.0170", "8.5149", "0.0983", "0.9679", "0.0033", "1.2436", "69"],
            ["FloED", "30.7459", "0.9570", "0.0263", "3.2203", "0.3365", "0.9689", "0.0034", "1.2124", "15"],
        ],
    },
    "Native-Best-24G / OR": {
        "kind": "OR",
        "conclusion": "结论：OR 的背景指标更像“保背景能力”，不能单独代表 object removal 成功率；这一页建议始终结合可视化一起读。",
        "rows": [
            ["ProPainter", "100.0000", "1.0000", "0.9582", "71"],
            ["FloED", "62.4027", "1.0000", "0.9560", "15"],
            ["DiffuEraser", "34.1224", "0.9819", "0.9578", "71"],
            ["MiniMax-Remover", "27.7824", "0.9396", "0.9718", "67"],
        ],
    },
    "Native-Best-Fix-240-432-24G / BR": {
        "kind": "BR",
        "conclusion": "结论：Fix-240-432 / BR 去掉了输入分辨率差异；ProPainter 仍最强，MiniMax 与 DiffuEraser 接近，FloED 维持稳定。",
        "rows": [
            ["ProPainter", "35.5167", "0.9813", "0.0141", "7.4734", "0.0818", "0.9703", "0.0045", "1.2362", "69"],
            ["MiniMax-Remover", "31.1977", "0.9612", "0.0190", "4.3716", "0.1019", "0.9688", "0.0010", "1.2299", "65"],
            ["DiffuEraser", "31.0785", "0.9637", "0.0172", "8.4298", "0.1048", "0.9681", "0.0097", "1.2397", "69"],
            ["FloED", "30.7687", "0.9572", "0.0269", "3.2117", "0.3313", "0.9688", "0.0059", "1.2118", "15"],
        ],
    },
    "Native-Best-Fix-240-432-24G / OR": {
        "kind": "OR",
        "conclusion": "结论：Fix-240-432 / OR 是低分辨率统一条件下的 OR 对照；背景指标仍可能饱和，因此建议配合 mask overlay 视频一起判断。",
        "rows": [
            ["DiffuEraser", "100.0000", "1.0000", "0.9531", "71"],
            ["ProPainter", "100.0000", "1.0000", "0.9582", "71"],
            ["FloED", "62.4027", "1.0000", "0.9573", "15"],
            ["MiniMax-Remover", "24.9133", "0.8921", "0.9681", "67"],
        ],
    },
    "Native-Best-Fix-nFrames-24G / BR": {
        "kind": "BR",
        "conclusion": "结论：Fix-nFrames / BR 去掉了长窗口差异；ProPainter 仍第一，MiniMax 领先于 DiffuEraser，FloED 的变化最小。",
        "rows": [
            ["ProPainter", "35.1758", "0.9783", "0.0157", "6.7386", "0.2269", "0.9723", "0.0123", "1.1812", "16"],
            ["MiniMax-Remover", "31.7330", "0.9636", "0.0213", "3.0584", "0.2590", "0.9714", "0.0394", "1.2370", "13"],
            ["DiffuEraser", "30.8651", "0.9600", "0.0178", "7.4467", "0.2819", "0.9714", "0.0292", "1.1952", "16"],
            ["FloED", "30.7840", "0.9566", "0.0276", "3.2158", "0.3281", "0.9691", "0.0050", "1.2195", "15"],
        ],
    },
    "Native-Best-Fix-nFrames-24G / OR": {
        "kind": "OR",
        "conclusion": "结论：Fix-nFrames / OR 统一了窗口长度，但背景指标读法不变；它更适合做同窗口对照，不适合单独宣告 removal 最优。",
        "rows": [
            ["ProPainter", "100.0000", "1.0000", "0.9620", "16"],
            ["FloED", "62.4027", "1.0000", "0.9571", "15"],
            ["DiffuEraser", "34.0971", "0.9813", "0.9622", "16"],
            ["MiniMax-Remover", "27.8303", "0.9409", "0.9716", "13"],
        ],
    },
    "Native-Best-Fix-nFrames-Fix-240-432-24G / BR": {
        "kind": "BR",
        "conclusion": "结论：这是 low-res source 下控制变量最强的 BR 表；适合做公平对照，但不是最终最重要的 full-res 主表。",
        "rows": [
            ["ProPainter", "35.1758", "0.9783", "0.0157", "6.7386", "0.2269", "0.9723", "0.0063", "1.1812", "16"],
            ["FloED", "30.6865", "0.9567", "0.0264", "3.2311", "0.3415", "0.9693", "0.0454", "1.2173", "15"],
            ["DiffuEraser", "30.6611", "0.9592", "0.0178", "7.4538", "0.2688", "0.9708", "0.0099", "1.1908", "16"],
            ["MiniMax-Remover", "29.7829", "0.9484", "0.0241", "3.5100", "0.3145", "0.9701", "0.0184", "1.2300", "13"],
        ],
    },
    "Native-Best-Fix-nFrames-Fix-240-432-24G / OR": {
        "kind": "OR",
        "conclusion": "结论：这是 low-res source 下控制变量最强的 OR 表；但 OR 最终仍应结合 mask overlay 视频一起解释。",
        "rows": [
            ["DiffuEraser", "100.0000", "1.0000", "0.9588", "16"],
            ["ProPainter", "100.0000", "1.0000", "0.9620", "16"],
            ["FloED", "62.4027", "1.0000", "0.9581", "15"],
            ["MiniMax-Remover", "24.7793", "0.8867", "0.9647", "13"],
        ],
    },
    "Native-Best-FullRes-24G / BR": {
        "kind": "BR",
        "conclusion": "结论：这是最接近“真实 high-res 输入 + 各方法 native-style 舒适区”的 BR 表，也是最重要的主表之一。",
        "rows": [
            ["ProPainter", "34.2067", "0.9782", "0.0162", "10.2557", "0.0806", "0.9748", "0.0033", "1.2537", "69"],
            ["MiniMax-Remover", "31.8785", "0.9670", "0.0173", "6.5358", "0.0850", "0.9746", "0.0044", "1.2391", "65"],
            ["DiffuEraser", "30.3474", "0.9613", "0.0173", "11.3108", "0.0964", "0.9720", "0.0050", "1.2606", "69"],
            ["FloED", "30.1085", "0.9536", "0.0276", "5.5943", "0.3402", "0.9744", "0.0191", "1.2454", "15"],
        ],
    },
    "Native-Best-Fix-240-432-FullRes-24G / BR": {
        "kind": "BR",
        "conclusion": "结论：这是最推荐的 BR 主表：真实 high-res 源图 + 统一 432x240 推理，既避免低分辨率图像上采样，又保留公平比较。",
        "rows": [
            ["ProPainter", "34.2067", "0.9782", "0.0162", "10.2557", "0.0806", "0.9748", "0.0038", "1.2537", "69"],
            ["MiniMax-Remover", "30.3931", "0.9579", "0.0188", "6.8619", "0.1040", "0.9734", "0.0083", "1.2516", "65"],
            ["DiffuEraser", "30.3591", "0.9612", "0.0173", "11.3890", "0.0952", "0.9720", "0.0048", "1.2592", "69"],
            ["FloED", "30.1468", "0.9538", "0.0271", "5.5549", "0.3190", "0.9751", "0.0060", "1.2352", "15"],
        ],
    },
    "Native-Best-Fix-nFrames-FullRes-24G / BR": {
        "kind": "BR",
        "conclusion": "结论：在 full-res 源图上只控制帧预算后，ProPainter 仍第一，MiniMax 与 DiffuEraser 更接近，可作为 supplementary control。",
        "rows": [
            ["ProPainter", "33.8962", "0.9751", "0.0177", "9.4691", "0.2375", "0.9763", "0.0172", "1.2041", "16"],
            ["MiniMax-Remover", "30.6123", "0.9584", "0.0201", "5.5313", "0.2547", "0.9764", "0.0091", "1.2561", "13"],
            ["DiffuEraser", "30.2752", "0.9584", "0.0175", "10.1341", "0.2677", "0.9743", "0.0231", "1.2080", "16"],
            ["FloED", "30.1928", "0.9547", "0.0265", "5.5868", "0.3199", "0.9749", "N/A", "1.2449", "15"],
        ],
    },
    "Native-Best-Fix-nFrames-Fix-240-432-FullRes-24G / BR": {
        "kind": "BR",
        "conclusion": "结论：这是 full-res 源图下同时控制分辨率与帧预算最严格的 BR 表，适合做最终 supplementary control。",
        "rows": [
            ["ProPainter", "33.8962", "0.9751", "0.0177", "9.4691", "0.2375", "0.9763", "0.0053", "1.2041", "16"],
            ["DiffuEraser", "30.3599", "0.9594", "0.0175", "10.1945", "0.2639", "0.9751", "0.0032", "1.2142", "16"],
            ["FloED", "30.2457", "0.9544", "0.0259", "5.5687", "0.3130", "0.9749", "0.0117", "1.2380", "15"],
            ["MiniMax-Remover", "29.1567", "0.9455", "0.0235", "5.9205", "0.3356", "0.9744", "0.0067", "1.2450", "13"],
        ],
    },
}


MAIN_TITLES = [
    "Native-Best-24G / BR",
    "Native-Best-24G / OR",
    "Native-Best-Fix-240-432-24G / BR",
    "Native-Best-Fix-240-432-24G / OR",
    "Native-Best-Fix-nFrames-24G / BR",
    "Native-Best-Fix-nFrames-24G / OR",
    "Native-Best-Fix-nFrames-Fix-240-432-24G / BR",
    "Native-Best-Fix-nFrames-Fix-240-432-24G / OR",
]

FULLRES_TITLES = [
    "Native-Best-FullRes-24G / BR",
    "Native-Best-Fix-240-432-FullRes-24G / BR",
    "Native-Best-Fix-nFrames-FullRes-24G / BR",
    "Native-Best-Fix-nFrames-Fix-240-432-FullRes-24G / BR",
]


def clear_slide(slide):
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape.element)


def clone_shapes(src_slide, dst_slide):
    for shape in src_slide.shapes:
        new_el = deepcopy(shape.element)
        dst_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")


def add_cloned_slide(prs, template_slide, position):
    blank = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank)
    clear_slide(new_slide)
    clone_shapes(template_slide, new_slide)
    sldIdLst = prs.slides._sldIdLst
    new_id = sldIdLst[-1]
    sldIdLst.remove(new_id)
    sldIdLst.insert(position, new_id)
    return prs.slides[position]


def set_shape_text(shape, text, size=18, bold=False):
    tf = shape.text_frame
    tf.clear()
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold


def update_page_numbers(prs):
    total = len(prs.slides)
    pat = re.compile(r"^\d+/\d+$")
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and pat.match(shape.text.strip()):
                shape.text = f"{i}/{total}"


def trim_table_rows(table, desired_rows):
    while len(table.rows) > desired_rows:
        table._tbl.remove(table._tbl.tr_lst[-1])


def set_table(table, headers, rows, font_size):
    trim_table_rows(table, len(rows) + 1)
    for c, value in enumerate(headers):
        table.cell(0, c).text = value
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    for row in table.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)


def update_result_slide(slide, title, seq_label, data):
    headers = {
        "BR": ["Method", "PSNR", "SSIM", "LPIPS", "Ewarp", "VFID", "TC", "AS", "IS", "Frames"],
        "OR": ["Method", "PSNR(bg)", "SSIM(bg)", "TC", "Frames"],
    }
    slide.shapes[3].text = f"固定 Benchmark 结果 {seq_label}：{title}"
    slide.shapes[4].text = data["conclusion"]
    table = next(shape.table for shape in slide.shapes if getattr(shape, "has_table", False))
    set_table(table, headers[data["kind"]], data["rows"], 12 if data["kind"] == "BR" else 14)


def update_summary_slide(slide):
    slide.shapes[3].text = "最重要的三个实验标准"
    body = (
        "主结论优先看这 3 张 BR 表：\n"
        "1. `Native-Best-FullRes-24G / BR`\n"
        "2. `Native-Best-Fix-240-432-FullRes-24G / BR`\n"
        "3. `Native-Best-Fix-240-432-24G / BR`\n\n"
        "为什么这三张最重要：\n"
        "• `Native-Best-FullRes-24G / BR`：真实 high-res 输入，方法可使用各自 native-style 舒适区。\n"
        "• `Native-Best-Fix-240-432-FullRes-24G / BR`：真实 high-res 输入，但统一到 432x240 推理，是最推荐主表。\n"
        "• `Native-Best-Fix-240-432-24G / BR`：低分辨率 source 的公平对照，用于验证排序是否依赖上采样。\n\n"
        "统一结论：从真实 high-res 输入下采样到低分辨率，明显比把 432x240 图像先上采样再推理更合理。"
    )
    slide.shapes[4].text = body


def update_slide9(slide):
    slide.shapes[4].text = "固定 Benchmark 的核心协议"
    slide.shapes[5].text = (
        "主结论优先看 3 张 BR 表：\n"
        "`Native-Best-FullRes-24G / BR`\n"
        "`Native-Best-Fix-240-432-FullRes-24G / BR`\n"
        "`Native-Best-Fix-240-432-24G / BR`\n"
        "原因：它们最接近“真实 high-res 输入 + native-style 舒适区 / 或统一 432x240 公平比较”，避免低分辨率图像先上采样再推理。"
    )


def update_later_conclusion_slides(prs):
    for slide in prs.slides:
        title = slide.shapes[3].text if len(slide.shapes) > 3 and hasattr(slide.shapes[3], "text") else ""
        if title == "正式实验前必须 Pre 的内容":
            slide.shapes[4].text = (
                "不需要特别展开的简单工程项：\n"
                "sbatch 封装 | cache 重定向 | 常规日志输出"
            )
            table = next(shape.table for shape in slide.shapes if getattr(shape, "has_table", False))
            table.cell(2, 1).text = (
                "先固定主 BR 表：优先看 `Native-Best-FullRes-24G / BR` 与 "
                "`Native-Best-Fix-240-432-FullRes-24G / BR`；若需要 low-res 对照，再看 "
                "`Native-Best-Fix-240-432-24G / BR`"
            )
        elif title == "建议的实验路线图":
            slide.shapes[4].text = (
                "1. 先锁定 BR 主 benchmark：`Native-Best-FullRes-24G / BR`、"
                "`Native-Best-Fix-240-432-FullRes-24G / BR`，并用 "
                "`Native-Best-Fix-240-432-24G / BR` 做 low-res 对照。\n"
                "2. 以当前 SFT DiffuEraser 作为 0 号 baseline，再做统一 benchmark。\n"
                "3. 在主表固定后，再推进 DPO smoke test 与正式训练。"
            )
        elif title == "结束页":
            slide.shapes[4].text = (
                "论文综述主线：方法谱系、数据集、setting、benchmark 已经完整闭环\n"
                "结果主线：固定 benchmark 的 12 张表已经更新进 PPT，可直接作为当前定量部分\n"
                "最重要的三张 BR 表：`Native-Best-FullRes-24G / BR`、"
                "`Native-Best-Fix-240-432-FullRes-24G / BR`、`Native-Best-Fix-240-432-24G / BR`\n"
                "我们自己的工作主线：SFT / DPO 继续围绕 DiffuEraser 主干推进，MiniMax 作为强 frozen 对照"
            )


def update_placeholder_slide(slide, title, body):
    slide.shapes[3].text = title
    slide.shapes[4].text = body


def main():
    prs = Presentation(PPT_PATH)

    template_br = prs.slides[9]
    template_summary = prs.slides[1]

    update_slide9(prs.slides[8])

    roman = ["I", "II", "III", "IV", "V", "VI", "VII", "VIII"]
    for idx, title in enumerate(MAIN_TITLES, start=9):
        update_result_slide(prs.slides[idx], title, roman[idx - 9], TABLES[title])

    # Repurpose old slides 18-20 into first three full-res BR tables.
    for slide_idx in [17, 18, 19]:
        clear_slide(prs.slides[slide_idx])
        clone_shapes(template_br, prs.slides[slide_idx])

    for idx, title in enumerate(FULLRES_TITLES[:3], start=17):
        seq = ["IX", "X", "XI"][idx - 17]
        update_result_slide(prs.slides[idx], title, seq, TABLES[title])

    # Add 4th full-res table after current slide 20, and summary after it.
    fullres4 = add_cloned_slide(prs, template_br, 20)
    update_result_slide(fullres4, FULLRES_TITLES[3], "XII", TABLES[FULLRES_TITLES[3]])

    summary_slide = add_cloned_slide(prs, template_summary, 21)
    update_summary_slide(summary_slide)

    # Add two placeholder slides at the end.
    ph1 = add_cloned_slide(prs, template_summary, len(prs.slides))
    update_placeholder_slide(
        ph1,
        "Diffueraser_vs_MiniMax 结果占位 I",
        "待插入表格：`Native-Best-24G / BR`\n\n"
        "内容：DiffuEraser 三个权重 + PCM / raw steps + MiniMax 6 / 12 / 50。",
    )
    ph2 = add_cloned_slide(prs, template_summary, len(prs.slides))
    update_placeholder_slide(
        ph2,
        "Diffueraser_vs_MiniMax 结果占位 II",
        "待插入表格：`Native-Best-FullRes-24G / BR`\n\n"
        "内容：DiffuEraser 三个权重 + PCM / raw steps + MiniMax 6 / 12 / 50。",
    )

    update_later_conclusion_slides(prs)
    update_page_numbers(prs)
    prs.save(PPT_PATH)


if __name__ == "__main__":
    main()
